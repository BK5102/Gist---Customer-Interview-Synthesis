# Text extraction for document types: PDF, PPTX, DOCX
from io import BytesIO


def extract_pdf(content: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(BytesIO(content))
    parts = []
    for page in reader.pages:
        text = page.extract_text() or ""
        if text.strip():
            parts.append(text.strip())
    return "\n\n".join(parts)


def extract_pptx(content: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(BytesIO(content))
    parts = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs if run.text.strip())
                    if line.strip():
                        slide_text.append(line.strip())
        if slide_text:
            parts.append("\n".join(slide_text))
    return "\n\n".join(parts)


def extract_docx(content: bytes) -> str:
    from docx import Document
    doc = Document(BytesIO(content))
    parts = [para.text.strip() for para in doc.paragraphs if para.text.strip()]
    return "\n\n".join(parts)


def extract_document(content: bytes, ext: str) -> str:
    """Dispatch to the correct extractor. Raise ValueError if output is empty."""
    extractors = {
        ".pdf": extract_pdf,
        ".pptx": extract_pptx,
        ".docx": extract_docx,
    }
    fn = extractors.get(ext)
    if fn is None:
        raise ValueError(f"Unsupported document extension: {ext}")
    text = fn(content)
    if not text.strip():
        raise ValueError(
            f"No readable text found in this {ext[1:].upper()} file. "
            "The file may be scanned, image-only, or password-protected."
        )
    return text
